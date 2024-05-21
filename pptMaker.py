from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.shapes.shapetree import GroupShapes
import json
import pprint

# Title (presentation title slide)
# Title and Content
# Section Header (sometimes called Segue)
# Two Content (side by side bullet textboxes)
# Comparison (same but additional title for each side by side content box)
# Title Only
# Blank
# Content with Caption
# Picture with Caption
TITLE_ONLY = 5
BLANK = 6
CONTENT_CAPTION = 7
PICTURE_CAPTION = 8

def make_slide(presentation) :
    slide_layout = presentation.slide_layouts[BLANK]
    return presentation.slides.add_slide(slide_layout)    

def make_image(shapes, location, size, img_path) :
    return shapes.add_picture(img_path, Cm(location[1]), Cm(location[0]), Cm(size[0]), Cm(size[1]))

def make_image(shapes, layer_position, layer_size, position, size ,img_path) :
    return shapes.add_picture(img_path, Cm(position[1]+layer_position[1]), Cm(position[0]+layer_position[0]), Cm(size[0]), Cm(size[1]))

def make_textbox(shapes, location, size) : 
    return shapes.add_textbox(Cm(location[0]), Cm(location[1]), Cm(size[0]), Cm(size[1])).text_frame

def make_textbox(shapes, layer_position, layer_size, position,size) :
    return shapes.add_textbox(Cm(position[0]+layer_position[0]), Cm(position[1]+layer_position[1]), Cm(size[0]), Cm(size[1])).text_frame


def make_text(text_frame, text, config) :
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(config.get('size'))
    run.font.name = config.get('font')
    
    color = config.get('color')
    run.font.color.rgb = RGBColor(color[0] , color[1] , color[2])

def make_layer(shapes, position, size) :
    shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(position[1]), Cm(position[0]), Cm(size[0]), Cm(size[1]))
    fill = shape.fill
    
    #transparent
    fill.background()
    line = shape.line
    line.color.rgb = RGBColor(0, 0, 0)
    
    shadow = shape.shadow
    shadow.inherit = False
    return shape
    
def read_json() :
    json_data = []
    with open('./data.json','r') as json_file:
            json_data = json.load(json_file)
    return json_data

def find_root_layer(data) :
    child_layer = []
    
    for layer_key in data :
        layer = data.get(layer_key)
        child_layer_list = layer.get('child_layer')
        
        for layer in child_layer_list :
            if layer not in child_layer :
                child_layer.append(layer)
        
    return data.keys() - child_layer
    
data = read_json()

root_layer = find_root_layer(data)

#work list
prs = Presentation()
slide = make_slide(prs)

def make_ppt(data, layer_list, root_position, group) :
    for key in layer_list :
        layer = data.get(key)
        
        size  = layer.get('size')
        position  = layer.get('position')

        position[0] = position[0] + root_position[0]
        position[1] = position[1] + root_position[1]
        
        shape_group = group.shapes.add_group_shape()
        
        make_layer(shape_group.shapes, position,  size )

        object_list = layer.get('object_list')
    
        for object in object_list :
            type = object.get('type')
            if type == 'text' :  
                obj_position = object.get('position')
                obj_size = object.get('size')
                text_list = object.get('text_list')
                text_frame = make_textbox(shape_group.shapes, position, size, obj_position, obj_size)
                
                for text_object in text_list :
                    text = text_object.get('text')
                    config = {
                        'font' : '나눔스퀘어 Bold',
                        'size' : 10,
                        'color' : (0,0,0)
                    }
                    make_text(text_frame, text , config)
            elif type == 'picture' :
                obj_position = object.get('position')
                obj_size = object.get('size')
                make_image(shape_group.shapes, position, size, obj_position, obj_size, 'test.jpg')
    
        child_list = layer.get('child_layer')
        make_ppt(data, child_list, position, shape_group)
    
make_ppt(data, root_layer,(0,0), slide)
    
def make_exection_message(message):
    return ""
    
prs.save("test.pptx")
